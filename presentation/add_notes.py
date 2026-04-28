"""
Inject speaker notes into the dissertation proposal deck.
25-minute talk: ~70 seconds per slide average.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from lxml import etree
import copy

PPTX_PATH = Path(__file__).resolve().parent / "dissertation_proposal_defense.pptx"

# Notes keyed by 1-based slide number
NOTES = {
    1: """Welcome everyone. Title: Bridging Domain Gaps in Visual Correspondence.
This proposal covers three chapters: Chapter 1 is completed published work on cross-modal aerial sensing, Chapter 2 is under review at ECCV and introduces directed coverage metrics, and Chapter 3 — which I'll spend the most time on — is my active research developing a principled data curation framework. The unifying theme across all three is a data-centric approach: treating training data composition as a rigorous, optimizable variable rather than an afterthought.""",

    2: """Dense visual correspondence is the core operation behind optical flow, semantic matching, 3D reconstruction, and navigation. Despite looking like different problems, they all share a single primitive: given a point in one image, find its counterpart in another. The challenge I'm addressing is that correspondence models trained on one distribution regularly fail when deployed on another. This happens due to sensing modality gaps, motion statistic differences, and fundamentally different supervision formats — dense flow versus sparse keypoints. My thesis is that a data-centric approach — rigorously analyzing and selecting training set composition — is the key to making progress on this.""",

    3: """The dissertation has three chapters that build on each other. Chapter 1 establishes the domain shift problem concretely through remote sensing — where the gap is literally between different sensing physics. Chapter 2 asks: can we measure which training datasets are likely to transfer, without running expensive downstream experiments? Chapter 3 turns that measurement into an action: given a pool of heterogeneous training pairs, which subset should we actually train on? The arrow here is important — the diagnostic tool from Chapter 2 becomes the utility function inside Chapter 3's selection algorithm.""",

    4: """My thesis in one sentence: domain shift failures in visual correspondence can be diagnosed and reduced through data-centric, motion-aware analysis and selection of training set composition. The three chapters each contribute one lever toward that claim. Chapter 1 establishes the problem domain concretely. Chapter 2 gives us the diagnostic — directed coverage tells us which training sets are likely to transfer and why. Chapter 3 turns that diagnostic into an action — given a heterogeneous pool, which pairs should we actually train on? The connecting thread is that motion structure is an explicit, measurable, actionable signal, not just a background variable.""",

    5: """Chapter 1 is completed, published work from CVPRW 2022 through 2024. The motivation is that remote sensing is the most extreme version of domain shift we'll encounter — EO captures reflected visible light, SAR measures radar backscatter, IR measures thermal emission, and magnetic maps encode geophysical structure. We organized the MAVOC and MAVIC challenge series using the UNICORN and MAGIC datasets. Three key findings: multi-modal training consistently improves single-modal inference; evaluation protocol matters — geographic disjoint splits revealed severe location overfitting that joint splits missed; and the map inference problem generalizes the same ideas to non-image domains through MINNIMAN.""",

    6: """MAVIC-T is our open-source translation benchmark — four translation directions across the MAGIC dataset. The headline finding is that SAR-to-RGB is substantially harder than SAR-to-EO because the sensing physics gap is larger and the model must hallucinate color structure. MINNIMAN takes the same cross-modal logic and applies it to navigation: given EO and hyperspectral inputs, predict magnetic anomaly maps that a GPS-denied vehicle could use as a reference. It's a proof-of-concept — the VQ-GAN codebook discretizes the map space and a teacher-student strategy handles changing input modalities. Both systems are published and complete.""",

    7: """Chapter 2 starts with a question: when a training dataset transfers well to a target, what property explains it? The standard tool is a symmetric distribution distance like MMD or FID. But symmetric distances have a fundamental problem — they conflate two qualitatively different failure modes into a single scalar. Failure mode A is under-coverage: the training set is missing motion patterns that the target requires. Failure mode B is excess mass: the training set contains lots of irrelevant patterns. These have different implications — missing support is usually much more harmful — and they require different remedies. MMD can't tell them apart.""",

    8: """Before I dive into Chapter 2, I want to make the motivation explicit. Chapter 1 demonstrated that even with one of the most carefully constructed multi-modal datasets in the field — MAGIC is one of the few open-source datasets that is spatially and temporally co-registered across SAR, RGB, IR, and hyperspectral — we still hit hard limits. Getting sensors to agree on the same patch of ground at the same moment requires precise flight geometry, stable scenes, and significant post-processing. And even then, we get no dense correspondence labels. This bottleneck is not unique to aerial sensing: dense optical flow requires calibrated stereo or LiDAR, semantic keypoints require human annotation at scale, and labels from one domain don't transfer to another. Chapter 2's starting point is: what if we sidestep the collection problem entirely and generate synthetic training pairs with full ground-truth geometry?""",

    9: """To test whether motion structure actually drives transfer, I built SDF-Fractal3D — an intentionally non-photorealistic synthetic pipeline. It renders paired views from implicit signed-distance geometry with procedural fractal textures. What makes it useful scientifically is that motion sampling and appearance generation are factorized: I can vary the zoom distribution or flip augmentation while holding appearance completely fixed. This lets me run clean controlled experiments. The surprising finding is that despite being visually abstract, models pretrained on SDF-Fractal3D transfer strongly zero-shot to real benchmarks — sometimes matching or beating photorealistic generators.""",

    10: """To compare datasets, I need a shared representation. For motion, I use Bag-of-Flow-Vectors: each correspondence becomes a 4D point — normalized source position x-hat, y-hat, and normalized displacement delta-x-hat, delta-y-hat. This is resolution-agnostic, captures both where and how motion occurs, and avoids grid discretization artifacts. Critically, because it's in Euclidean space with physical scale, I can define a meaningful epsilon-coverage radius. For appearance, I use DINOv2 patch embeddings projected onto the unit hypersphere. Together these two descriptor spaces let me compute asymmetric coverage in both directions.""",

    11: """Directed coverage gives me two numbers for each training-target pair. E ⊆ε T coverage measures: what fraction of target BFV descriptors have a nearby neighbor in the training set? Low E ⊆ε T means the training data is missing motion modes the benchmark requires — this is the primary failure signal. T ⊆ε E measures: what fraction of training descriptors are near the target? Low T ⊆ε E means the training set wastes capacity on irrelevant patterns. By computing both directions in both motion and appearance spaces, I get four interpretable predictors. These are lightweight — no training required, just nearest-neighbor queries on descriptor clouds.""",

    12: """The motion tuning table shows the controlled experiment directly — only the motion sampler varies, appearance is held fixed. The key results: large zoom augmentation adds six to seven PCK@5% points on TSS and four points on KITTI-15. Flip mismatch — training with horizontal flips on a benchmark that doesn't have them — costs five to eight points. These are large, systematic effects explained entirely by directed coverage changes, not by any appearance difference. This validates BFV as an actionable intervention signal, not just a post-hoc diagnostic.""",

    13: """PCK — Percentage of Correct Keypoints — is the standard metric across semantic matching and template benchmarks. A prediction counts as correct if it falls within alpha times the max image dimension of the ground-truth location. In Chapter 2 we use alpha equals five percent throughout. In Chapter 3 the ablations use the default alpha of ten percent. Higher is better; no fine-tuning on any target benchmark.""",

    14: """The transferability estimator takes the four directed coverage predictors as input and uses ridge regression to predict within-context performance residuals. I evaluate under three increasingly strict held-out protocols: held-out target benchmark, held-out training dataset, and both simultaneously. The key results: the directed 4-predictor model — Flow times 2 plus Appearance times 2 — achieves 0.69 to 0.70 pairwise ranking accuracy across all three protocols. Chance is 0.50. Symmetric MMD actually falls below chance on appearance, which motivates the asymmetric formulation. Motion terms are the strongest signal; appearance adds consistent secondary value.""",

    15: """Chapter 2 delivers three things that feed directly into Chapter 3. First, a motion descriptor — BFV — that captures distributional fingerprints without rigid quantization. Second, a diagnostic framework that separates two distinct failure modes that symmetric metrics conflate. Third, a compact transferability estimator that ranks training datasets without downstream retraining. The bridge to Chapter 3 is conceptually clean: E ⊆ε T coverage — which measures missing target support at the dataset level — becomes the per-pair scoring function for greedy subset selection at the pair level. The same signal that explains transfer now drives curation.""",

    16: """Chapter 3 addresses a practical problem that arises as soon as you have multiple correspondence sources. My pool has three: PointOdyssey with about 4.4 million pairs of dense synthetic flow, SPair-71k with 53 thousand pairs of sparse semantic keypoints, and PF-PASCAL with under 3 thousand pairs. That's a 98.7% to 1.2% to 0.07% composition. Naive pooling fails for three reasons: PointOdyssey's consecutive frames generate massive near-duplicate redundancy; dense sources dominate coverage scores purely by volume; and there's no principled mechanism for allocating budget across sources with incomparable supervision formats.""",

    17: """My initial approach is ClusterCov: cluster the target benchmark's BFV space into K motion modes, then greedily select training pairs that cover uncovered clusters. Step 4 — supervision-density normalization — is critical. Without it, a PointOdyssey pair with 5000 flow vectors gets the same score as an SPair pair with 12 keypoints, even though the PointOdyssey pair has already explained most of its motion in the first few correspondences. Dividing by the number of valid supervision points turns raw coverage gain into per-observation efficiency, which is what lets the selector fairly compare across sources.""",

    18: """These figures show the Pareto tradeoff. The left panel makes the key point visually: no single source Pareto-dominates across both geometric and semantic targets. PointOdyssey-only is very strong on KITTI but weak on PF-PASCAL and TSS. PF-PASCAL-only flips that. The multi-target ClusterCov selector sits in the upper-right, outperforming pooled random and approaching the hand-tuned baseline — without any manual source balancing. The right panel shows this across all five benchmarks at convergence. The claim I'm making is not that ClusterCov beats every single-source oracle; it's Pareto dominance across the full target set.""",

    19: """This is the core unsolved problem that motivates the GIST framework. The density mismatch figure shows it concretely. SPair and PF-PASCAL pairs have roughly 12 keypoints each after rasterization into BFV space — you can see the sparse, cluster-like BFV cloud. PointOdyssey pairs have 128-plus flow vectors — the BFV cloud is dense and fills the space. When you feed both into the same coverage selector, PointOdyssey dominates by sheer count even after normalization, because its BFV cloud saturates coverage faster. And critically — the bottom-right panel shows that this same density mismatch causes the MAE latent space to cluster by dataset identity rather than by motion structure, which is the central challenge for the stretch goal.""",

    20: """The GIST framework formalizes the selection problem as maximizing a joint objective: coverage utility g(S) plus lambda times min-pairwise diversity. The coverage term is monotone non-negative submodular, which gives us a (1 minus 1/e) greedy guarantee on its own. Adding the diversity term breaks submodularity because diversity is monotone non-increasing. GIST solves this by sweeping a geometric grid of diversity thresholds and running constrained greedy at each, achieving a (one-half minus eta) approximation. We make three correspondence-specific adaptations: using directed BFV coverage as the utility instead of margin uncertainty; a source-partitioned distance metric that makes cross-source pairs always diverse; and a supervision-density-derived coverage radius that removes per-target manual tuning.""",

    21: """This table is one of the most striking empirical results. The selector operates on motion descriptors alone — no class labels, no appearance features, no source-identity flags. Yet when targeting PF-WILLOW it draws 78% from SPair — a 66-times enrichment over SPair's 1.2% pool fraction. When targeting KITTI it flips to 53-57% PointOdyssey. TSS draws 81% from SPair. This routing matches exactly what a domain expert would assemble by hand, but it emerges purely from the geometry of motion distribution overlap. This is independent evidence that BFV coverage encodes cross-domain compatibility structure that goes beyond its geometric framing.""",

    22: """The masked correspondence autoencoder is a stretch goal, but I think it's genuinely interesting as a research problem in its own right. The question is whether we can learn latents that are invariant to observation sparsity — so that a pair with 12 keypoints and a pair with 5000 flow vectors land in the same region of latent space if they describe similar motion. The architecture uses a ViT encoder with 75% masking on patchified RGB and observed flow. The problem I hit is that the latent space clusters by dataset identity rather than motion structure — which is shown directly on the preceding slide. My proposed fix is replacing the linear patch embedding with Perceiver IO, which treats keypoints within a patch as a set and produces a fixed-dimension token regardless of density. Whether this resolves the clustering or just relocates it to the ViT encoder is an open question.""",

    23: """The validation plan has three parts corresponding to the three chapters. Chapter 1 is already validated — published results. For Chapter 2 the key tests are the three held-out ranking protocols and the controlled SDF motion interventions. For Chapter 3: does joint normalized ClusterCov beat pooled random at matched budget, does it produce sensible source routing, and does the GIST lambda generalize from calibration to deployment benchmarks? If lambda generalizes — practical offline curation tool. If not — Pareto frontier generator. Either result is publishable.""",

    24: """Chapters 1 and 2 are done. I'm currently running Chapter 3 curation experiments and source ablations in parallel. The target is to have the main Chapter 3 experiments wrapped by May 30th, the latent model investigation closed by June 27th, an advisor draft by July 25th, and the committee-review submission out by August 1st. The defense is planned for August 21st. The timeline is tight but tractable given the job offer driving the deadline.""",

    32: """Thank you. Happy to take questions on any aspect — the directed coverage framework, the GIST formalization, the preliminary curation results, or the MAE latent model direction.""",
}

def add_notes(slide, text):
    """Inject speaker notes XML into a slide."""
    # Get or create notes slide
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = text
    # Set font size
    from pptx.util import Pt
    for run in p.runs:
        run.font.size = Pt(11)

prs = Presentation(str(PPTX_PATH))

for slide_num, note_text in NOTES.items():
    idx = slide_num - 1
    if idx < len(prs.slides):
        add_notes(prs.slides[idx], note_text.strip())
        print(f"  ✓ Slide {slide_num} notes added")

prs.save(str(PPTX_PATH))
print(f"\n✅  Saved with speaker notes: {PPTX_PATH}")

# ── Hide backup/appendix slides ────────────────────────────────
# After the 4b+4c merge AND thesis slide insertion the deck is 31 slides (0-indexed 0–30).
# Slide order (0-indexed):
#  0  Title
#  1  Motivation
#  2  Task gallery (2b)
#  3  Chapter overview
#  4  Thesis statement  (NEW)
#  5  Ch1 overview
#  6  Ch1 combined MAVIC-T + MINNIMAN  (merged 4b+4c)
#  7  Ch2 why symmetric fails
#  8  Directed coverage framework
#  9  BFV descriptor
# 10  Dataset visual gallery (8b)           ← HIDE
# 11  SDF-Fractal3D
# 12  Motion tuning table
# 13  PCK primer
# 14  Transferability estimator
# 15  Ch2 summary
# 16  Transfer grid (11b)                   ← HIDE
# 17  Ch3 opening
# 18  ClusterCov
# 19  Converged PCK
# 20  Convergence curves (13b)              ← HIDE
# 21  Source ablations
# 22  Cross-source budget allocation
# 23  GIST framework
# 24  Source routing table
# 25  MAE stretch goal
# 26  Validation plan
# 27  Timeline (treat as last presented slide)
# 28  Contributions at a glance             ← HIDE
# 29  Anticipated questions                 ← HIDE
# 30  Thank You / Closing

SLIDES_TO_HIDE = [21, 29, 30]  # 0-indexed

prs2 = Presentation(str(PPTX_PATH))
for idx in SLIDES_TO_HIDE:
    if idx < len(prs2.slides):
        prs2.slides[idx]._element.attrib['show'] = '0'
        print(f"  🙈 Slide {idx+1} hidden")
prs2.save(str(PPTX_PATH))
print(f"✅  Hidden {len(SLIDES_TO_HIDE)} backup slides")
